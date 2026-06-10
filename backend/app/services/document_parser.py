"""文档解析管线 - PDF/DOCX/PPTX/MD/TXT"""

from __future__ import annotations

from pathlib import Path
from typing import List

from loguru import logger


class DocumentParser:
    """文档解析器"""

    @staticmethod
    async def parse(file_path: str, file_type: str) -> str:
        """解析文档为纯文本"""
        path = Path(file_path)
        if not path.exists():
            raise FileNotFoundError(f"文件不存在: {file_path}")

        if file_type == "pdf":
            return await DocumentParser._parse_pdf(path)
        elif file_type == "docx":
            return await DocumentParser._parse_docx(path)
        elif file_type == "pptx":
            return await DocumentParser._parse_pptx(path)
        elif file_type == "md":
            return await DocumentParser._parse_md(path)
        elif file_type == "txt":
            return await DocumentParser._parse_txt(path)
        else:
            raise ValueError(f"不支持的文件类型: {file_type}")

    @staticmethod
    async def _parse_pdf(path: Path) -> str:
        """解析 PDF (纯 Python，无需编译)"""
        try:
            from pdfminer.high_level import extract_text
            text = extract_text(str(path))
            return text.strip()
        except ImportError:
            logger.warning("pdfminer 未安装")
            return f"[PDF 文件: {path.name}]"

    @staticmethod
    async def _parse_docx(path: Path) -> str:
        """解析 DOCX"""
        try:
            from docx import Document
            doc = Document(str(path))
            return "\n".join(p.text for p in doc.paragraphs if p.text.strip())
        except ImportError:
            logger.warning("python-docx 未安装")
            return f"[DOCX 文件: {path.name}]"

    @staticmethod
    async def _parse_pptx(path: Path) -> str:
        """解析 PPTX"""
        try:
            from pptx import Presentation
            prs = Presentation(str(path))
            texts = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        texts.append(shape.text)
            return "\n".join(texts)
        except ImportError:
            logger.warning("python-pptx 未安装")
            return f"[PPTX 文件: {path.name}]"

    @staticmethod
    async def _parse_md(path: Path) -> str:
        """解析 Markdown"""
        return path.read_text(encoding="utf-8")

    @staticmethod
    async def _parse_txt(path: Path) -> str:
        """解析纯文本"""
        return path.read_text(encoding="utf-8")
